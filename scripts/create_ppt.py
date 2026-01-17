from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Solarized Light Palette
    COLORS = {
        'base3': RGBColor(253, 246, 227),  # #fdf6e3 Background
        'base2': RGBColor(238, 232, 213),  # #eee8d5
        'base1': RGBColor(147, 161, 161),  # #93a1a1
        'base00': RGBColor(101, 123, 131), # #657b83 Text
        'base01': RGBColor(88, 110, 117),  # #586e75
        'yellow': RGBColor(181, 137, 0),   # #b58900
        'orange': RGBColor(203, 75, 22),   # #cb4b16 Title
        'red':    RGBColor(220, 50, 47),   # #dc322f
        'magenta': RGBColor(211, 54, 130), # #d33682
        'violet': RGBColor(108, 113, 196), # #6c71c4
        'blue':   RGBColor(38, 139, 210),  # #268bd2
        'cyan':   RGBColor(42, 161, 152),  # #2aa198
        'green':  RGBColor(133, 153, 0),   # #859900
    }

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['base3']

    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        set_background(slide)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text
        
        # Style Title
        title_tf = title.text_frame
        title_tf.paragraphs[0].font.color.rgb = COLORS['orange']
        title_tf.paragraphs[0].font.size = Pt(54)
        title_tf.paragraphs[0].font.bold = True
        
        # Style Subtitle
        subtitle_tf = subtitle.text_frame
        p = subtitle_tf.paragraphs[0]
        p.font.color.rgb = COLORS['base01']
        p.font.size = Pt(28)
        
        return slide

    def add_content_slide(title_text, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        
        # Style Title
        title_tf = title.text_frame
        title_tf.paragraphs[0].font.color.rgb = COLORS['yellow']
        title_tf.paragraphs[0].font.size = Pt(40)
        title_tf.paragraphs[0].font.bold = True
        
        # Content
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()  # Clear default empty paragraph
        
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.color.rgb = COLORS['base00']
            p.font.size = Pt(24)
            p.space_after = Pt(20)
            
    def add_two_column_slide(title_text, col1_title, col1_items, col2_title, col2_items):
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Using Title and Content as base
        set_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS['yellow']
        
        # Manually create two text boxes
        left = Inches(0.5)
        top = Inches(2.0)
        width = Inches(4.5)
        height = Inches(5.0)
        
        # Column 1
        txBox1 = slide.shapes.add_textbox(left, top, width, height)
        tf1 = txBox1.text_frame
        
        p = tf1.add_paragraph()
        p.text = col1_title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['blue']
        p.space_after = Pt(20)
        
        for item in col1_items:
            p = tf1.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['base00']
            p.level = 0
            
        # Column 2
        txBox2 = slide.shapes.add_textbox(left + Inches(5.0), top, width, height)
        tf2 = txBox2.text_frame
        
        p = tf2.add_paragraph()
        p.text = col2_title
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['blue']
        p.space_after = Pt(20)
        
        for item in col2_items:
            p = tf2.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['base00']
            p.level = 0

    def add_stats_slide(title_text, stats):
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
        set_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLORS['yellow']
        
        # Grid of stats
        start_x = Inches(0.5)
        start_y = Inches(2.5)
        box_width = Inches(2.2)
        box_height = Inches(1.5)
        gap = Inches(0.25)
        
        for i, (value, label) in enumerate(stats):
            left = start_x + (i * (box_width + gap))
            
            # Draw box background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, start_y, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLORS['base2']
            shape.line.color.rgb = COLORS['base1']
            
            # Value
            txt_box = slide.shapes.add_textbox(left, start_y + Inches(0.2), box_width, Inches(0.8))
            p = txt_box.text_frame.paragraphs[0]
            p.text = value
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(32)
            p.font.color.rgb = COLORS['magenta']
            
            # Label
            label_box = slide.shapes.add_textbox(left, start_y + Inches(0.8), box_width, Inches(0.5))
            p = label_box.text_frame.paragraphs[0]
            p.text = label
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['base01']

    # --- SLIDES ---

    # 1. Title
    add_title_slide(
        "Car Retrieval System",
        "An End-to-End Deep Learning Approach for\nIndonesian Vehicle Detection and Classification\n\nMachine Learning Engineering Challenge\nJanuary 2026"
    )

    # 2. Problem Statement
    add_two_column_slide(
        "Problem Statement",
        "Challenge", [
            "Develop Object Detection Model for Cars",
            "Classify Indonesian Car Types",
            "Integrate both into unified pipeline",
            "Process real-world traffic video"
        ],
        "Requirements", [
            "Detect multiple car instances",
            "8 car type categories",
            "CNN-based feature extraction",
            "Deep Learning Framework (PyTorch)"
        ]
    )

    # 3. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide)
    
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['yellow']
    
    # Text representation of flow
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    p = flow_box.text_frame.paragraphs[0]
    p.text = "Input Image/Video  →  YOLOv8 Detector  →  Crop Regions  →  ResNet50 Classifier  →  Car Types"
    p.font.size = Pt(18)
    p.font.name = "Courier New"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = COLORS['base00']
    
    # Details
    add_two_column_slide(
        "System Details", 
        "Detection Stage", ["YOLOv8-nano (COCO pretrained)", "Filters vehicle classes", "Returns BBoxes + Confidence"],
        "Classification Stage", ["ResNet50 backbone (ImageNet)", "Custom classification head", "8 output classes"]
    )
    # Correcting the flow: remove the extra slide, merge into one if possible. 
    # For simplicity, let's keep the details on the separate slide created by add_two_column_slide
    # But wait, add_two_column_slide creates NEW slide. 
    # Let's just create a custom single slide for Arch.
    
    # 4. Dataset
    add_stats_slide("Dataset Overview", [
        ("17,171", "Total Images"),
        ("8", "Car Types"),
        ("13,542", "Train Samples"),
        ("1,756", "Test Samples")
    ])
    # Add classes text
    slide = prs.slides[-1]
    classes_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1))
    p = classes_box.text_frame.add_paragraph()
    p.text = "Crossover • Hatchback • MPV • Offroad • Pickup • Sedan • Truck • Van"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['base01']

    # 5. Classifier Arch
    add_two_column_slide(
        "Classifier Architecture",
        "ResNet50 Backbone", [
            "Pre-trained on ImageNet",
            "50 convolutional layers",
            "Skip connections",
            "24.5M parameters"
        ],
        "Custom Head", [
            "Global Average Pooling",
            "Dropout (p=0.5)",
            "FC: 2048 → 512 → 8",
            "Softmax output"
        ]
    )

    # 6. Results Overview
    add_stats_slide("Results Overview", [
        ("86.33%", "Test Accuracy"),
        ("86.53%", "Precision"),
        ("86.34%", "F1-Score"),
        ("13.6", "Video FPS")
    ])

    # 7. Per-Class Table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide)
    
    title = slide.shapes.title
    title.text = "Per-Class Performance"
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['yellow']
    
    table_placeholder = slide.shapes.add_table(9, 4, Inches(1), Inches(2.0), Inches(8), Inches(4.5))
    table = table_placeholder.table
    
    # Headers
    headers = ["Class", "Precision", "Recall", "F1-Score"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['base2']
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = COLORS['orange']

    # Data
    data = [
        ("Crossover", 0.85, 0.81, 0.83),
        ("Hatchback", 0.88, 0.83, 0.85),
        ("MPV", 0.80, 0.90, 0.85),
        ("Offroad", 0.86, 0.85, 0.86),
        ("Pickup", 0.93, 0.89, 0.91),
        ("Sedan", 0.87, 0.87, 0.87),
        ("Truck", 0.92, 0.88, 0.90),
        ("Van", 0.90, 0.90, 0.90)
    ]
    
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.color.rgb = COLORS['base00']
            
            # Highlight high scores
            if isinstance(val, float) and val >= 0.90:
                cell.text_frame.paragraphs[0].font.color.rgb = COLORS['green']
                cell.text_frame.paragraphs[0].font.bold = True

    # 8. Video Inference
    add_stats_slide("Video Inference Results", [
        ("2,960", "Frames Processed"),
        ("73.7ms", "Avg Inference"),
        ("14,107", "Car Detections"),
        ("197s", "Video Duration")
    ])
    
    slide = prs.slides[-1]
    dist_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1))
    p = dist_box.text_frame.add_paragraph()
    p.text = "MPV: 34.7% • Sedan: 19.5% • Hatchback: 16.1% • Offroad: 12.0%"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['base01']

    # 9. Key Findings
    add_two_column_slide(
        "Key Findings",
        "Insights", [
            "Pickup & Truck: High accuracy (>90%) due to distinct shapes",
            "MPV: High recall but confused with Crossovers",
            "Real-time capable at 13.6 FPS"
        ],
        "Tech Stack", [
            "PyTorch",
            "Ultralytics YOLOv8",
            "torchvision",
            "OpenCV",
            "scikit-learn"
        ]
    )

    # 10. Conclusion
    add_two_column_slide(
        "Conclusion",
        "Achievements", [
            "86.33% Test Accuracy",
            "All 8 car types handled",
            "End-to-End Pipeline",
            "Video & Image Support"
        ],
        "Future Work", [
            "Vision Transformers (ViT)",
            "Dataset Expansion",
            "Edge Deployment (TensorRT)",
            "Vehicle Tracking (DeepSORT)"
        ]
    )

    # Save
    prs.save('presentation.pptx')
    print("Presentation saved as presentation.pptx")

if __name__ == "__main__":
    create_presentation()
